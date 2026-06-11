from fastapi import FastAPI, File, UploadFile, Form
from fastapi.responses import FileResponse
from fastapi.staticfiles import StaticFiles
from fastapi.middleware.cors import CORSMiddleware
from PIL import Image, ImageOps
from pptx import Presentation
from pptx.util import Inches, Pt, Cm
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
import os
import io
import json
import numpy as np
from scipy.signal import find_peaks
from scipy.ndimage import gaussian_filter1d
from typing import Optional

app = FastAPI()

app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)

os.makedirs("static", exist_ok=True)
app.mount("/static", StaticFiles(directory="static"), name="static")

@app.get("/")
def read_root():
    return FileResponse("static/index.html")

def detect_bands(marker_image, expected_count):
    # Auto-contrast to enhance faint bands
    gray_image = ImageOps.autocontrast(marker_image.convert('L'))
    data = np.array(gray_image)
    row_avgs = np.mean(data, axis=1)
    signal = 255 - row_avgs
    
    # Smooth signal to reduce noise
    signal = gaussian_filter1d(signal, sigma=2)
    
    # Subtract low-frequency background (handles uneven illumination/gradient)
    bg_sigma = max(15, len(signal) // 10)
    background = gaussian_filter1d(signal, sigma=bg_sigma)
    signal = signal - background
    
    # Suppress edges
    edge = max(2, int(len(signal) * 0.015))
    signal[:edge] = np.min(signal)
    signal[-edge:] = np.min(signal)
    
    distance = max(2, len(signal) // (expected_count * 2.5))
    sig_range = np.max(signal) - np.min(signal)
    prom_threshold = max(0.1, sig_range * 0.01)
    
    peaks, properties = find_peaks(signal, distance=distance, prominence=prom_threshold)
    prominences = properties['prominences']
    
    # Get the top expected_count peaks based on prominence
    sorted_peak_indices = np.argsort(prominences)[::-1]
    top_peaks = peaks[sorted_peak_indices[:expected_count]]
    top_peaks.sort()
    return top_peaks

MARKERS = {
    "commercial PAGE": ["170", "130", "93", "72", "53", "42", "30", "23", "14", "10"],
    "15% PAGE": ["180", "75", "60", "45", "35", "25", "15", "10"],
    "10% PAGE": ["180", "140", "100", "75", "60", "45", "35"]
}

def process_page_image(image_file, lanes_str, scale_factor, temp_path):
    lanes_data = json.loads(lanes_str)
    if not lanes_data:
        return None, None, None, None
        
    img_data = image_file.file.read()
    original_image = Image.open(io.BytesIO(img_data)).convert('RGB')
    
    global_y1 = min(l['y1'] for l in lanes_data)
    global_y2 = max(l['y2'] for l in lanes_data)
    
    cropped_images = []
    marker_image = None
    
    for lane in lanes_data:
        x1 = int(lane['x1'] / scale_factor)
        x2 = int(lane['x2'] / scale_factor)
        gy1 = int(global_y1 / scale_factor)
        gy2 = int(global_y2 / scale_factor)
        
        cropped = original_image.crop((x1, gy1, x2, gy2))
        cropped_images.append(cropped)
        
        if lane['name'].strip().upper() in ['M', 'MARKER']:
            marker_image = cropped
            
    total_width = sum(img.width for img in cropped_images)
    max_height = int((global_y2 - global_y1) / scale_factor)
    stitched = Image.new('RGB', (total_width, max_height), color='white')
    
    x_offset = 0
    lane_centers = []
    for img in cropped_images:
        stitched.paste(img, (x_offset, 0))
        lane_centers.append(x_offset + img.width / 2)
        x_offset += img.width
        
    stitched.save(temp_path)
    return temp_path, total_width, max_height, lane_centers, marker_image, lanes_data

@app.post("/generate_ppt")
async def generate_ppt(
    image_1: UploadFile = File(...),
    lanes_1: str = Form(...),
    marker_name_1: str = Form(...),
    ignore_top_n_1: int = Form(0),
    scale_factor_1: float = Form(1.0),
    
    image_2: Optional[UploadFile] = File(None),
    lanes_2: Optional[str] = Form(None),
    marker_name_2: Optional[str] = Form(None),
    ignore_top_n_2: int = Form(0),
    scale_factor_2: float = Form(1.0),
    
    conditions: str = Form(...),
    yields: str = Form(...)
):
    try:
        cond_data = json.loads(conditions)
        yields_data = json.loads(yields)
        
        prs = Presentation()
        # 10.0 x 7.5 inches (4:3)
        prs.slide_width = Inches(10)
        prs.slide_height = Inches(7.5)
        
        slide_layout = prs.slide_layouts[6] # Blank
        slide = prs.slides.add_slide(slide_layout)
        
        # --- TITLE ---
        protein = cond_data.get('protein', '').strip()
        txBox_title = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9.0), Inches(0.4))
        tf_title = txBox_title.text_frame
        tf_title.clear()
        p = tf_title.paragraphs[0]
        p.text = f"Purification of : {protein}"
        p.font.bold = True
        p.font.size = Pt(18)
        p.font.name = "微軟正黑體"
        
        # --- CONDITIONS ---
        plasmid = cond_data.get('plasmid', '').strip()
        ecoli = cond_data.get('ecoli', '').strip()
        hek = cond_data.get('hek', '').strip()
        cho = cond_data.get('cho', '').strip()
        vol_l = cond_data.get('vol_l', '').strip()
        
        exp_parts = []
        if ecoli: exp_parts.append(f"E.coli: {ecoli}")
        if hek: exp_parts.append(f"HEK293: {hek}")
        if cho: exp_parts.append(f"CHO: {cho}")
        expr_str = ", ".join(exp_parts)
        if vol_l: expr_str += f" ({vol_l} L)"
        
        binding_type = cond_data.get('binding_type', '').strip()
        binding_detail = cond_data.get('binding_detail', '').strip()
        bind_str = binding_type
        if binding_detail: bind_str += f" / {binding_detail}"
        
        wash = cond_data.get('wash', '').strip()
        elute = cond_data.get('elute', '').strip()
        dialysis = cond_data.get('dialysis', '').strip()

        info_lines = [
            ("Protein/ Antibody: ", protein),
            ("Plasmid: ", plasmid),
            ("Expression: ", expr_str),
            ("Binding: ", bind_str),
            ("Wash: ", wash),
            ("Elution: ", elute),
            ("Dialysis buffer: ", dialysis)
        ]

        txBox_info = slide.shapes.add_textbox(Inches(0.5), Inches(0.9), Inches(9.0), Inches(2.0))
        tf_info = txBox_info.text_frame
        tf_info.clear()
        
        for i, (label, val) in enumerate(info_lines):
            p = tf_info.paragraphs[0] if i == 0 else tf_info.add_paragraph()
            p.space_after = Pt(2)
            run1 = p.add_run()
            run1.text = label
            run1.font.bold = True
            run1.font.size = Pt(12)
            run1.font.name = "微軟正黑體"
            
            run2 = p.add_run()
            run2.text = val
            run2.font.size = Pt(12)
            run2.font.name = "微軟正黑體"
            
        # --- YIELD BLOCK (BOTTOM HORIZONTAL) ---
        purification_vol_L = None
        try:
            if vol_l: purification_vol_L = float(vol_l)
        except ValueError:
            pass
            
        yield_y = 6.4
        txBox_bot = slide.shapes.add_textbox(Inches(0.8), Inches(yield_y), Inches(8.4), Inches(1.0))
        tf_bot = txBox_bot.text_frame
        tf_bot.clear()
        tf_bot.margin_top = 0
        tf_bot.margin_left = 0
        
        valid_yield_idx = 0
        for row in yields_data:
            s_name = row.get('name', '').strip()
            s_conc = row.get('conc', '').strip()
            s_vol = row.get('vol', '').strip()
            
            if s_conc and s_vol:
                try:
                    conc_val = float(s_conc)
                    vol_val_ml = float(s_vol)
                    total_mg = conc_val * vol_val_ml
                    yield_str = "0"
                    if purification_vol_L and purification_vol_L > 0:
                        yield_str = f"{(total_mg / purification_vol_L):.2f}"
                        
                    line_str = f"{s_name}產量: 濃度: {conc_val:.2f} mg/mL, 體積: {vol_val_ml:.2f} mL, 總量: {total_mg:.2f} mg, 產率: {yield_str} mg/L"
                    
                    p = tf_bot.paragraphs[0] if valid_yield_idx == 0 else tf_bot.add_paragraph()
                    p.space_after = Pt(4)
                    p.text = line_str
                    p.font.size = Pt(12)
                    p.font.name = "微軟正黑體"
                    
                    valid_yield_idx += 1
                except ValueError:
                    pass

        # Helper to draw PAGE image with annotations
        def draw_page(img_file, lanes_str, scale_fac, marker_name, ignore_n, title, title_x, title_y, img_x, img_y, temp_path):
            txBox = slide.shapes.add_textbox(Inches(title_x), Inches(title_y), Inches(3.0), Inches(0.4))
            tf = txBox.text_frame
            tf.clear()
            p = tf.paragraphs[0]
            p.text = title
            p.font.bold = True
            p.font.size = Pt(12)
            p.font.name = "微軟正黑體"
            p.alignment = PP_ALIGN.CENTER

            if not img_file or not lanes_str:
                # Add placeholder
                shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(img_x), Inches(img_y), Inches(2.5), Inches(3.4))
                shape.fill.solid()
                shape.fill.fore_color.rgb = RGBColor(255, 255, 255)
                shape.line.color.rgb = RGBColor(153, 153, 153)
                shape.line.width = Pt(1)
                tf = shape.text_frame
                tf.text = f"(請在此處貼入 {title} 圖)"
                tf.paragraphs[0].font.size = Pt(10)
                tf.paragraphs[0].font.color.rgb = RGBColor(136, 136, 136)
                tf.paragraphs[0].font.name = "微軟正黑體"
                tf.paragraphs[0].alignment = PP_ALIGN.CENTER
                return

            # Process image
            res = process_page_image(img_file, lanes_str, scale_fac, temp_path)
            if not res or not res[0]: return
            path, total_width, max_height, lane_centers, marker_image, lanes_data = res
            
            img_height = Inches(3.2)
            pic = slide.shapes.add_picture(path, Inches(img_x), Inches(img_y), height=img_height)
            img_width = pic.width 
            ratio = img_width / total_width

            sample_counter = 1
            for i, lane in enumerate(lanes_data):
                is_marker = lane['name'].strip().upper() in ['M', 'MARKER']
                display_name = 'M' if is_marker else str(sample_counter)
                if not is_marker: sample_counter += 1
                
                center_x = pic.left + lane_centers[i] * ratio
                top = pic.top - Cm(0.6)
                width = Cm(1.5)
                txBox2 = slide.shapes.add_textbox(center_x - width/2, top, width, Cm(0.6))
                tf2 = txBox2.text_frame
                tf2.clear()
                tf2.margin_bottom = 0
                tf2.margin_top = 0
                tf2.margin_left = 0
                tf2.margin_right = 0
                tf2.vertical_anchor = MSO_ANCHOR.BOTTOM
                p2 = tf2.paragraphs[0]
                p2.text = display_name
                p2.font.size = Pt(12)
                p2.font.name = "微軟正黑體"
                p2.alignment = PP_ALIGN.CENTER

            # Add Legend to the right of the image
            legend_x = pic.left + pic.width + Cm(0.2)
            legend_y = pic.top
            legend_txBox = slide.shapes.add_textbox(legend_x, legend_y, Cm(4.0), pic.height)
            legend_tf = legend_txBox.text_frame
            legend_tf.clear()
            legend_tf.margin_top = 0
            legend_tf.margin_left = 0
            
            non_marker_lanes = [l for l in lanes_data if l['name'].strip().upper() not in ['M', 'MARKER']]
            for i, lane in enumerate(non_marker_lanes):
                p_leg = legend_tf.paragraphs[0] if i == 0 else legend_tf.add_paragraph()
                p_leg.text = f"{i + 1}: {lane['name']}"
                p_leg.font.size = Pt(10)
                p_leg.font.name = "微軟正黑體"

            if marker_image:
                marker_values = MARKERS.get(marker_name, MARKERS["commercial PAGE"])
                expected_count = len(marker_values) + ignore_n
                try:
                    peaks_y = detect_bands(marker_image, expected_count)
                    if ignore_n > 0 and len(peaks_y) > ignore_n:
                        peaks_y = peaks_y[ignore_n:]
                        
                    for i, val in enumerate(marker_values):
                        if i < len(peaks_y):
                            y_pixel = peaks_y[i]
                            y_ppt = pic.top + (y_pixel / max_height) * img_height
                            label_width = Cm(1.5)
                            label_left = pic.left - label_width - Cm(0.1)
                            label_top = y_ppt - Cm(0.3) 
                            txBox3 = slide.shapes.add_textbox(label_left, label_top, label_width, Cm(0.6))
                            tf3 = txBox3.text_frame
                            tf3.clear()
                            tf3.margin_bottom = 0
                            tf3.margin_top = 0
                            tf3.margin_left = 0
                            tf3.margin_right = 0
                            tf3.vertical_anchor = MSO_ANCHOR.MIDDLE
                            p3 = tf3.paragraphs[0]
                            p3.text = str(val)
                            p3.font.size = Pt(10)
                            p3.font.name = "微軟正黑體"
                            p3.alignment = PP_ALIGN.RIGHT
                except Exception as e:
                    print(f"Band detection error: {e}")

        # 8. Draw Purify and Dialysis
        draw_page(image_1, lanes_1, scale_factor_1, marker_name_1, ignore_top_n_1, "純化PAGE", 0.8, 2.7, 0.8, 3.0, "temp_purify.png")
        draw_page(image_2, lanes_2, scale_factor_2, marker_name_2, ignore_top_n_2, "透析PAGE", 6.0, 2.7, 6.0, 3.0, "temp_dialysis.png")

        output_path = "SDS_PAGE_Result.pptx"
        prs.save(output_path)
        
        for p in ["temp_purify.png", "temp_dialysis.png"]:
            if os.path.exists(p):
                os.remove(p)
            
        final_name = f"蛋白純化報告(4比3)_{protein if protein else '未命名'}.pptx"
        return FileResponse(output_path, media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation", filename=final_name)

    except Exception as e:
        import traceback
        traceback.print_exc()
        return {"error": str(e)}
